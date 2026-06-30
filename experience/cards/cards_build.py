# -*- coding: utf-8 -*-
"""Build Particleworks Experience 2026 speaker cards: flat PNGs + editable square PPTX."""
import os, numpy as np, cv2
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageOps, ImageChops

DL   = r"C:\Users\e4user\Downloads"
ROOT = r"C:\Users\e4user\Documents\ParticleworksProjects\pw-website\experience"
OUT  = os.path.join(ROOT, "cards")
PROC = os.path.join(DL, "_proc"); os.makedirs(PROC, exist_ok=True); os.makedirs(OUT, exist_ok=True)
FONT = os.path.join(DL, "OfficinaSans-Bold.otf")
BG   = os.path.join(DL, "backgroun_new.png")
LOGO = os.path.join(DL, "experience_logo.png")
SPK  = os.path.join(ROOT, "speakers")

S=1200; M=90
WHITE=(255,255,255); GREEN=(52,210,124); SUB=(255,224,216); SUBLOC=(255,217,207)
C_IN=(220,96,78); C_OUT=(150,28,30)          # reddish avatar gradient
def f(sz): return ImageFont.truetype(FONT, int(sz))

# ---- layout template, taken from the user-edited example slides 1-3 (px @120dpi) ----
P2X=1.6667                                    # pt -> px at 120dpi  (px = pt*120/72)
CAL_POS=(81,86); CAL_W=78; CAL_H=84
DATE_POS=(168,86); DATE_PX=42                 # 25.2pt
CAT_POS=(168,140); CAT_PX=26                  # 15.6pt
TITLE_X,TITLE_Y,TITLE_W=90,238,1020
TITLE_MAXPX,TITLE_MINPX,TITLE_MAXH=90,58,468  # 54pt, shrink long titles to fit
AV_X,AV_Y,AV_D=106,702,209
COMP_X,COMP_Y=345,734; PRES_X,PRES_Y=345,794; CP_PX=60   # 36pt company + presenter
LOGO_X,LOGO_Y,LOGO_W=90,1009,348
VEN_RIGHT=1110; VEN1_Y,VEN1_PX=1020,53; VEN2_Y,VEN2_PX=1072,33   # 32pt / 20pt
SIM_X,SIM_Y,SIM_W,SIM_H=701,590,409,344

# ---------- asset prep ----------
def seg_cutout(path,out):
    img=cv2.imread(path); H0,W0=img.shape[:2]
    sc=max(1,700//max(H0,W0));
    work=cv2.resize(img,(W0*sc,H0*sc),interpolation=cv2.INTER_CUBIC) if sc>1 else img.copy()
    h,w=work.shape[:2]
    pil=Image.fromarray(cv2.cvtColor(work,cv2.COLOR_BGR2RGB)); flood=pil.copy(); sent=(0,255,1)
    for s in [(0,0),(w-1,0),(0,h-1),(w-1,h-1),(w//2,0),(0,h//2),(w-1,h//2),(w//4,0),(3*w//4,0)]:
        ImageDraw.floodfill(flood,s,sent,thresh=40)
    defbg=np.all(np.asarray(flood)==sent,axis=-1)
    gm=np.full((h,w),2,np.uint8); gm[defbg]=0
    yy,xx=np.ogrid[:h,:w]; gm[((xx-w/2)/(w*0.32))**2+((yy-h*0.46)/(h*0.42))**2<=1]=3
    cv2.grabCut(work,gm,None,np.zeros((1,65),np.float64),np.zeros((1,65),np.float64),7,cv2.GC_INIT_WITH_MASK)
    m=np.where((gm==1)|(gm==3),255,0).astype('uint8')
    m=cv2.morphologyEx(m,cv2.MORPH_CLOSE,np.ones((11,11),np.uint8))
    m=cv2.morphologyEx(m,cv2.MORPH_OPEN,np.ones((5,5),np.uint8))
    n,lab,st,_=cv2.connectedComponentsWithStats(m,8)
    if n>1: m=np.where(lab==1+np.argmax(st[1:,cv2.CC_STAT_AREA]),255,0).astype('uint8')
    m=cv2.GaussianBlur(m,(0,0),2.0)
    m=cv2.resize(m,(W0,H0),interpolation=cv2.INTER_AREA)
    base=Image.open(path).convert("RGB"); rgba=base.convert("RGBA"); rgba.putalpha(Image.fromarray(m))
    rgba.save(out)

def radial(dia,ci,co):
    yy,xx=np.mgrid[0:dia,0:dia].astype(float)
    r=np.clip(np.sqrt((xx-dia/2)**2+(yy-dia/2)**2)/(dia/2),0,1)[...,None]
    return Image.fromarray((np.array(ci)*(1-r)+np.array(co)*r).astype('uint8'),"RGB")

def make_avatar(kind,val,dia=320):
    out=Image.new("RGBA",(dia,dia),(0,0,0,0))
    cmask=Image.new("L",(dia,dia),0); ImageDraw.Draw(cmask).ellipse([0,0,dia-1,dia-1],fill=255)
    out.paste(radial(dia,C_IN,C_OUT),(0,0),cmask)
    if kind=="photo":
        person=ImageOps.fit(Image.open(val).convert("RGBA"),(dia,dia),Image.LANCZOS,centering=(0.5,0.42))
        person.putalpha(ImageChops.multiply(person.split()[3],cmask))
        out.alpha_composite(person)
    else:
        d=ImageDraw.Draw(out); ft=f(int(dia*0.4)); bb=ft.getbbox(val)
        d.text((dia/2-ft.getlength(val)/2, dia/2-(bb[3]-bb[1])/2-bb[1]),val,font=ft,fill=WHITE)
    ImageDraw.Draw(out).ellipse([3,3,dia-4,dia-4],outline=WHITE,width=int(dia*0.03))
    return out

def make_sim(path,out,w=924,h=744):
    img=ImageOps.fit(Image.open(path).convert("RGB"),(w,h),Image.LANCZOS)
    mask=Image.new("L",(w,h),0); ImageDraw.Draw(mask).rounded_rectangle([0,0,w,h],radius=40,fill=255)
    r=Image.new("RGBA",(w,h),(0,0,0,0)); r.paste(img,(0,0)); r.putalpha(mask)
    ImageDraw.Draw(r).rounded_rectangle([2,2,w-3,h-3],radius=40,outline=WHITE,width=7)
    r.save(out)

def square_bg(out):
    im=Image.open(BG).convert("RGB"); w,h=im.size; sc=S/w
    im=im.resize((S,int(h*sc)),Image.LANCZOS); top=int((im.size[1]-S)*0.36); im=im.crop((0,top,S,top+S))
    scrim=Image.new("L",(S,S),0); sd=ImageDraw.Draw(scrim)
    for x in range(S): sd.line([(x,0),(x,S)],fill=max(45,int(120-55*(x/S))))
    im=Image.composite(Image.new("RGB",(S,S),(60,8,10)),im,scrim)
    sc2=Image.new("L",(S,S),0); d2=ImageDraw.Draw(sc2)
    for y in range(S): d2.line([(0,y),(S,y)],fill=int(max(0,(y-820)/(S-820))*110) if y>820 else 0)
    Image.composite(Image.new("RGB",(S,S),(40,6,8)),im,sc2).save(out)

def cal_icon(out,s=70,col=WHITE):
    im=Image.new("RGBA",(s+8,s+12),(0,0,0,0)); d=ImageDraw.Draw(im)
    d.rounded_rectangle([2,10,s+2,s+6],radius=7,outline=col,width=4); d.rectangle([2,10,s+2,26],fill=col)
    d.line([16,3,16,16],fill=col,width=4); d.line([s-12,3,s-12,16],fill=col,width=4)
    for r in range(3):
        for c in range(3): d.ellipse([14+c*16,33+r*9,19+c*16,38+r*9],fill=col)
    im.save(out)

# ---------- text helpers ----------
def fit(text,max_w,start,lo=22):
    sz=start
    while sz>lo:
        if f(sz).getlength(text)<=max_w: return f(sz)
        sz-=2
    return f(lo)
def wrap(text,ft,max_w):
    out=[]; cur=""
    for w in text.split():
        t=(cur+" "+w).strip()
        if ft.getlength(t)<=max_w: cur=t
        else:
            if cur: out.append(cur)
            cur=w
    if cur: out.append(cur)
    return out

# ---------- shared text fitting ----------
def fit_title(text,maxh=TITLE_MAXH,minpx=TITLE_MINPX):
    px=TITLE_MAXPX
    while px>=minpx:
        ft=f(px); lines=wrap(text,ft,TITLE_W); lh=int(px*1.16)
        if len(lines)*lh<=maxh: return ft,lines,lh
        px-=3
    ft=f(minpx); return ft,wrap(text,ft,TITLE_W),int(minpx*1.16)
def title_box(c):
    # when a sim panel sits lower-right (top at SIM_Y), keep the title above it
    if c.get("sim"): return dict(maxh=SIM_Y-TITLE_Y-14, minpx=46)
    return dict(maxh=TITLE_MAXH, minpx=TITLE_MINPX)
def comp_w(sim): return (SIM_X-COMP_X-20) if sim else (VEN_RIGHT-COMP_X)

# ---------- flat PNG renderer ----------
def render_png(c):
    im=Image.open(os.path.join(PROC,"bg_square.png")).convert("RGBA"); d=ImageDraw.Draw(im)
    d.rectangle([30,30,S-31,S-31],outline=WHITE,width=3); d.rectangle([48,48,S-49,S-49],outline=GREEN,width=5)
    cal=Image.open(os.path.join(PROC,"cal.png")).resize((CAL_W,CAL_H),Image.LANCZOS); im.alpha_composite(cal,CAL_POS)
    d=ImageDraw.Draw(im)
    d.text(DATE_POS,c["date"],font=f(DATE_PX),fill=WHITE); d.text(CAT_POS,c["cat"],font=f(CAT_PX),fill=SUB)
    sim=c.get("sim")
    if c.get("title2"):
        # two stacked titles (keynote): part 1 on top (white), part 2 below (green), divider between
        tf,lines,lh=fit_title(c["title"],maxh=212,minpx=48); y=TITLE_Y
        for ln in lines: d.text((TITLE_X,y),ln,font=tf,fill=WHITE); y+=lh
        y+=14; d.line([(TITLE_X,y),(TITLE_X+130,y)],fill=GREEN,width=5); y+=30
        t2=fit(c["title2"],TITLE_W,60,lo=40); lh2=int(t2.size*1.16)
        for ln in wrap(c["title2"],t2,TITLE_W): d.text((TITLE_X,y),ln,font=t2,fill=GREEN); y+=lh2
    else:
        tf,lines,lh=fit_title(c["title"],**title_box(c)); y=TITLE_Y
        for ln in lines: d.text((TITLE_X,y),ln,font=tf,fill=WHITE); y+=lh
    if sim:
        sp=Image.open(os.path.join(PROC,sim)).convert("RGBA").resize((SIM_W,SIM_H),Image.LANCZOS)
        sh=Image.new("RGBA",(S,S),(0,0,0,0)); ImageDraw.Draw(sh).rounded_rectangle([SIM_X,SIM_Y+8,SIM_X+SIM_W,SIM_Y+SIM_H+8],radius=40,fill=(0,0,0,95))
        im.alpha_composite(sh.filter(ImageFilter.GaussianBlur(12))); im.alpha_composite(sp,(SIM_X,SIM_Y))
    av=Image.open(os.path.join(PROC,c["avatar"])).convert("RGBA").resize((AV_D,AV_D),Image.LANCZOS)
    sh=Image.new("RGBA",(S,S),(0,0,0,0)); ImageDraw.Draw(sh).ellipse([AV_X,AV_Y+8,AV_X+AV_D,AV_Y+AV_D+8],fill=(0,0,0,90))
    im.alpha_composite(sh.filter(ImageFilter.GaussianBlur(10))); im.alpha_composite(av,(AV_X,AV_Y))
    d=ImageDraw.Draw(im); cw=comp_w(sim)
    cf=fit(c["company"],cw,CP_PX); pf=fit(c["speakers"],cw,CP_PX)
    d.text((COMP_X,COMP_Y),c["company"],font=cf,fill=WHITE); d.text((PRES_X,PRES_Y),c["speakers"],font=pf,fill=GREEN)
    lg=Image.open(LOGO).convert("RGBA"); lg=lg.resize((LOGO_W,int(lg.size[1]*LOGO_W/lg.size[0])),Image.LANCZOS)
    im.alpha_composite(lg,(LOGO_X,LOGO_Y))
    d.text((VEN_RIGHT-f(VEN1_PX).getlength(L1),VEN1_Y),L1,font=f(VEN1_PX),fill=WHITE)
    d.text((VEN_RIGHT-f(VEN2_PX).getlength(L2),VEN2_Y),L2,font=f(VEN2_PX),fill=SUBLOC)
    im.convert("RGB").save(os.path.join(OUT,c["file"]),"PNG"); print("png",c["file"])

# ---------- general "program is out" announcement card ----------
GEN_TAGS=["Surface Treatment","CFD-DEM","E-Motor Cooling","Reciprocating Pumps","Hydropower",
          "Thermal Simulation","Oil Aeration","Gear Lubrication","Undercarriage"]
def draw_tags(im,tags,x0,y0,maxw,fs=40,gap=18,vgap=20,padx=32,pady=15):
    # draw on a separate RGBA layer so the translucent pill fill blends instead of overwriting
    ov=Image.new("RGBA",im.size,(0,0,0,0)); d=ImageDraw.Draw(ov)
    ft=f(fs); ph=fs+pady*2; x=x0; y=y0
    for t in tags:
        pw=ft.getlength(t)+padx*2
        if x+pw>x0+maxw: x=x0; y+=ph+vgap
        d.rounded_rectangle([x,y,x+pw,y+ph],radius=ph//2,fill=(255,255,255,38),outline=GREEN,width=3)
        d.text((x+padx,y+ph/2),t,font=ft,fill=WHITE,anchor="lm")
        x+=pw+gap
    im.alpha_composite(ov)
    return y+ph
def render_general(fname="card-00-program.png"):
    im=Image.open(os.path.join(PROC,"bg_square.png")).convert("RGBA"); d=ImageDraw.Draw(im)
    d.rectangle([30,30,S-31,S-31],outline=WHITE,width=3); d.rectangle([48,48,S-49,S-49],outline=GREEN,width=5)
    cal=Image.open(os.path.join(PROC,"cal.png")).resize((CAL_W,CAL_H),Image.LANCZOS); im.alpha_composite(cal,CAL_POS)
    d=ImageDraw.Draw(im)
    d.text(DATE_POS,"October 6–7, 2026",font=f(DATE_PX),fill=WHITE)
    d.text(CAT_POS,"Conference programme",font=f(CAT_PX),fill=SUB)
    d.text((90,250),"THE PROGRAM",font=f(118),fill=WHITE)
    d.text((90,392),"IS OUT",font=f(118),fill=GREEN)
    d.text((92,556),"Particleworks Experience 2026 — 10 talks across two days",font=f(40),fill=SUBLOC)
    yend=draw_tags(im,GEN_TAGS,92,616,1016,fs=37,gap=16,vgap=16,padx=28,pady=13); d=ImageDraw.Draw(im)
    d.text((92,min(yend+24,952)),"Full agenda online · register free",font=f(42),fill=WHITE)
    lg=Image.open(LOGO).convert("RGBA"); lg=lg.resize((LOGO_W,int(lg.size[1]*LOGO_W/lg.size[0])),Image.LANCZOS)
    im.alpha_composite(lg,(LOGO_X,LOGO_Y))
    d.text((VEN_RIGHT-f(VEN1_PX).getlength(L1),VEN1_Y),L1,font=f(VEN1_PX),fill=WHITE)
    d.text((VEN_RIGHT-f(VEN2_PX).getlength(L2),VEN2_Y),L2,font=f(VEN2_PX),fill=SUBLOC)
    im.convert("RGB").save(os.path.join(OUT,fname),"PNG"); print("png",fname)

# ---------- Prometech keynote: two title+speaker segments ----------
def render_keynote(c):
    im=Image.open(os.path.join(PROC,"bg_square.png")).convert("RGBA"); d=ImageDraw.Draw(im)
    d.rectangle([30,30,S-31,S-31],outline=WHITE,width=3); d.rectangle([48,48,S-49,S-49],outline=GREEN,width=5)
    cal=Image.open(os.path.join(PROC,"cal.png")).resize((CAL_W,CAL_H),Image.LANCZOS); im.alpha_composite(cal,CAL_POS)
    d=ImageDraw.Draw(im)
    d.text(DATE_POS,c["date"],font=f(DATE_PX),fill=WHITE); d.text(CAT_POS,c["cat"],font=f(CAT_PX),fill=SUB)
    segs=[("PART 1 · SOFTWARE RELEASE",c["title"],"av_IM.png","Issei Masaie","General Manager · Prometech Software",WHITE),
          ("PART 2 · APPLICATIONS",c["title2"],"av_iori-saigo.png","Iori Saigo","Application Engineer · Prometech Software",GREEN)]
    starts=[202,560]
    for (kick,title,av,name,role,col),sy in zip(segs,starts):
        d.text((92,sy),kick,font=f(27),fill=SUB)
        tf,lines,lh=fit_title(title,maxh=150,minpx=44); ty=sy+40
        for ln in lines: d.text((92,ty),ln,font=tf,fill=col); ty+=lh
        avd=92; ay=ty+16
        avi=Image.open(os.path.join(PROC,av)).convert("RGBA").resize((avd,avd),Image.LANCZOS); im.alpha_composite(avi,(92,ay))
        tx=92+avd+20
        d.text((tx,ay+6),name,font=fit(name,360,40,lo=30),fill=WHITE)
        d.text((tx,ay+52),role,font=f(25),fill=SUBLOC)
    d.line([(92,524),(232,524)],fill=GREEN,width=5)
    lg=Image.open(LOGO).convert("RGBA"); lg=lg.resize((LOGO_W,int(lg.size[1]*LOGO_W/lg.size[0])),Image.LANCZOS)
    im.alpha_composite(lg,(LOGO_X,LOGO_Y))
    d.text((VEN_RIGHT-f(VEN1_PX).getlength(L1),VEN1_Y),L1,font=f(VEN1_PX),fill=WHITE)
    d.text((VEN_RIGHT-f(VEN2_PX).getlength(L2),VEN2_Y),L2,font=f(VEN2_PX),fill=SUBLOC)
    im.convert("RGB").save(os.path.join(OUT,c["file"]),"PNG"); print("png",c["file"],"(keynote)")

# ---------- carousel slide 2: speaker profiles + company tags ----------
SPEAKERS=[("iori-saigo","Iori Saigo","Prometech Software"),
 ("alpcan-guray","Alpcan Güray","MTU Aero Engines"),
 ("michelangelo-raimondo","Michelangelo Raimondo","UNIMORE"),
 ("leonardo-lanciotti","Leonardo Lanciotti","R&D CFD"),
 ("jean-decaix","Jean Decaix","HES-SO Valais"),
 ("lijun-cao","Lijun Cao","SKF"),
 ("lukas-hafner","Dr. Lukas Hafner","deepfluid"),
 ("leonardo-tiberi","Leonardo Tiberi","Track One"),
 ("naohiro-fujita","Naohiro Fujita","Univance"),
 ("rene-kockisch","René Kockisch","IAV")]
def render_speakers(fname="card-00b-speakers.png"):
    im=Image.open(os.path.join(PROC,"bg_square.png")).convert("RGBA"); d=ImageDraw.Draw(im)
    d.rectangle([30,30,S-31,S-31],outline=WHITE,width=3); d.rectangle([48,48,S-49,S-49],outline=GREEN,width=5)
    cal=Image.open(os.path.join(PROC,"cal.png")).resize((CAL_W,CAL_H),Image.LANCZOS); im.alpha_composite(cal,CAL_POS)
    d=ImageDraw.Draw(im)
    d.text(DATE_POS,"October 6–7, 2026",font=f(DATE_PX),fill=WHITE)
    d.text(CAT_POS,"Conference speakers",font=f(CAT_PX),fill=SUB)
    d.text((90,228),"MEET THE SPEAKERS",font=f(80),fill=WHITE)
    d.text((92,326),"10 industrial & academic talks · 6–7 Oct, Modena",font=f(34),fill=SUBLOC)
    AVD=96; cols=[92,612]; rows=[400,520,640,760,880]
    ov=Image.new("RGBA",im.size,(0,0,0,0)); od=ImageDraw.Draw(ov)
    for idx,(key,name,comp) in enumerate(SPEAKERS):
        cx=cols[idx//5]; cy=rows[idx%5]
        av=Image.open(os.path.join(PROC,"av_"+key+".png")).convert("RGBA").resize((AVD,AVD),Image.LANCZOS)
        im.alpha_composite(av,(cx,cy)); tx=cx+AVD+22
        d.text((tx,cy+10),name,font=fit(name,388,40,lo=27),fill=WHITE)
        cf=f(26); ph=26+16
        od.rounded_rectangle([tx,cy+56,tx+cf.getlength(comp)+36,cy+56+ph],radius=ph//2,fill=(255,255,255,38),outline=GREEN,width=2)
        od.text((tx+18,cy+56+ph/2),comp,font=cf,fill=WHITE,anchor="lm")
    im.alpha_composite(ov)
    lg=Image.open(LOGO).convert("RGBA"); lg=lg.resize((LOGO_W,int(lg.size[1]*LOGO_W/lg.size[0])),Image.LANCZOS)
    im.alpha_composite(lg,(LOGO_X,LOGO_Y))
    d.text((VEN_RIGHT-f(VEN1_PX).getlength(L1),VEN1_Y),L1,font=f(VEN1_PX),fill=WHITE)
    d.text((VEN_RIGHT-f(VEN2_PX).getlength(L2),VEN2_Y),L2,font=f(VEN2_PX),fill=SUBLOC)
    im.convert("RGB").save(os.path.join(OUT,fname),"PNG"); print("png",fname)

# ---------- editable PPTX ----------
def build_pptx(cards):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    def IN(px): return Inches(px/120.0)
    def PT(px): return Pt(px*0.6)
    rgb=lambda t: RGBColor(*t)
    prs=Presentation(); prs.slide_width=IN(S); prs.slide_height=IN(S)
    blank=prs.slide_layouts[6]
    def tb(sl,x,y,w,h,text,px,color,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
        b=sl.shapes.add_textbox(IN(x),IN(y),IN(w),IN(h)); tf=b.text_frame; tf.word_wrap=True
        tf.vertical_anchor=anchor; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text
        r.font.size=PT(px); r.font.bold=True; r.font.name="Officina Sans Bold"; r.font.color.rgb=rgb(color)
        return b
    for c in cards:
        sl=prs.slides.add_slide(blank)
        sl.shapes.add_picture(os.path.join(PROC,"bg_square.png"),0,0,IN(S),IN(S))
        for inset,col,wd in [(30,WHITE,2),(48,GREEN,3)]:
            fr=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(inset),IN(inset),IN(S-2*inset),IN(S-2*inset))
            fr.fill.background(); fr.line.color.rgb=rgb(col); fr.line.width=Pt(wd); fr.shadow.inherit=False
        sl.shapes.add_picture(os.path.join(PROC,"cal.png"),IN(CAL_POS[0]),IN(CAL_POS[1]),IN(CAL_W),IN(CAL_H))
        tb(sl,DATE_POS[0],DATE_POS[1],520,52,c["date"],DATE_PX,WHITE)
        tb(sl,CAT_POS[0],CAT_POS[1],520,34,c["cat"],CAT_PX,SUB)
        if c["file"]=="card-01-prometech.png":
            segs=[("PART 1 · SOFTWARE RELEASE",c["title"],"av_IM.png","Issei Masaie","General Manager · Prometech Software",WHITE),
                  ("PART 2 · APPLICATIONS",c["title2"],"av_iori-saigo.png","Iori Saigo","Application Engineer · Prometech Software",GREEN)]
            for (kick,title,av,name,role,col),sy in zip(segs,[202,560]):
                tb(sl,92,sy,900,34,kick,27,SUB)
                tf2,lines,lh=fit_title(title,maxh=150,minpx=44)
                tb(sl,TITLE_X,sy+40,TITLE_W,len(lines)*lh+10,title,tf2.size,col)
                ay=sy+40+len(lines)*lh+16
                sl.shapes.add_picture(os.path.join(PROC,av),IN(92),IN(ay),IN(92),IN(92))
                tx=92+92+20
                tb(sl,tx,ay+6,360,52,name,fit(name,360,40,lo=30).size,WHITE)
                tb(sl,tx,ay+52,880,40,role,25,SUBLOC)
            dv=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,IN(92),IN(524),IN(140),IN(6))
            dv.fill.solid(); dv.fill.fore_color.rgb=rgb(GREEN); dv.line.fill.background(); dv.shadow.inherit=False
            lg=Image.open(LOGO); lh0=int(lg.size[1]*LOGO_W/lg.size[0])
            sl.shapes.add_picture(LOGO,IN(LOGO_X),IN(LOGO_Y),IN(LOGO_W),IN(lh0))
            tb(sl,VEN_RIGHT-560,VEN1_Y,560,65,L1,VEN1_PX,WHITE,align=PP_ALIGN.RIGHT)
            tb(sl,VEN_RIGHT-560,VEN2_Y,560,40,L2,VEN2_PX,SUBLOC,align=PP_ALIGN.RIGHT)
            continue
        sim=c.get("sim")
        if c.get("title2"):
            tpx=fit_title(c["title"],maxh=212,minpx=48)[0].size
            tb(sl,TITLE_X,TITLE_Y,TITLE_W,236,c["title"],tpx,WHITE)
            t2px=fit(c["title2"],TITLE_W,60,lo=40).size
            tb(sl,TITLE_X,TITLE_Y+250,TITLE_W,180,c["title2"],t2px,GREEN)
        else:
            tpx=fit_title(c["title"],**title_box(c))[0].size
            tb(sl,TITLE_X,TITLE_Y,TITLE_W,TITLE_MAXH,c["title"],tpx,WHITE)
        if sim: sl.shapes.add_picture(os.path.join(PROC,sim),IN(SIM_X),IN(SIM_Y),IN(SIM_W),IN(SIM_H))
        sl.shapes.add_picture(os.path.join(PROC,c["avatar"]),IN(AV_X),IN(AV_Y),IN(AV_D),IN(AV_D))
        cw=comp_w(sim); cpx=fit(c["company"],cw,CP_PX).size; ppx=fit(c["speakers"],cw,CP_PX).size
        tb(sl,COMP_X,COMP_Y,cw,73,c["company"],cpx,WHITE)
        tb(sl,PRES_X,PRES_Y,cw,73,c["speakers"],ppx,GREEN)
        lg=Image.open(LOGO); lh=int(lg.size[1]*LOGO_W/lg.size[0])
        sl.shapes.add_picture(LOGO,IN(LOGO_X),IN(LOGO_Y),IN(LOGO_W),IN(lh))
        tb(sl,VEN_RIGHT-560,VEN1_Y,560,65,L1,VEN1_PX,WHITE,align=PP_ALIGN.RIGHT)
        tb(sl,VEN_RIGHT-560,VEN2_Y,560,40,L2,VEN2_PX,SUBLOC,align=PP_ALIGN.RIGHT)
    out=os.path.join(OUT,"particleworks-experience-2026-cards.pptx"); prs.save(out); print("pptx ->",out)

# ---------- data ----------
PHOTOS={"iori-saigo":"iori-saigo.jpg","alpcan-guray":"alpcan-guray.jpg","michelangelo-raimondo":"michelangelo-raimondo.jpg",
        "leonardo-lanciotti":"leonardo-lanciotti.jpg","jean-decaix":"jean-decaix.jpg","lukas-hafner":"lukas-hafner.jpg",
        "naohiro-fujita":"naohiro-fujita.jpg","lijun-cao":"lijun-cao.jpg","leonardo-tiberi":"leonardo-tiberi.jpg",
        "rene-kockisch":"rene-kockisch.jpg"}
DATE="October 7, 2026"; L1="Modena, Italy"; L2="BPER FORUM Monzani"
cards=[
 dict(file="card-01-prometech.png",date=DATE,cat="Developer keynote",company="PROMETECH SOFTWARE",speakers="Issei Masaie · Iori Saigo",
      title="What's New in Particleworks 9.0 & Granuleworks 4.0",
      title2="Application Examples & Case Studies",avatar="av_iori-saigo.png"),
 dict(file="card-02-mtu.png",date=DATE,cat="Industrial speaker",company="MTU AERO ENGINES",speakers="Alpcan Güray",
      title="Simulation of Shot Peening: A CFD-DEM Coupled Case Study",avatar="av_alpcan-guray.png"),
 dict(file="card-03-unimore-stator.png",date=DATE,cat="Academic speaker",company="UNIVERSITY OF MODENA AND REGGIO EMILIA",
      speakers="Michelangelo Raimondo",title="An Integrated Simulation Approach for Stator Oil Jacket and Jet Cooling Systems",avatar="av_michelangelo-raimondo.png"),
 dict(file="card-04-rdcfd-pump.png",date=DATE,cat="Industrial speaker",company="R&D CFD",speakers="Leonardo Lanciotti",
      title="CHT Analysis of a Reciprocating Pump",avatar="av_leonardo-lanciotti.png"),
 dict(file="card-05-hesso-pelton.png",date=DATE,cat="Academic speaker",company="HES-SO VALAIS//WALLIS",speakers="Jean Decaix",
      title="Moving Particle Simulation of Eroded Pelton Runners",avatar="av_jean-decaix.png"),
 dict(file="card-06-skf.png",date=DATE,cat="Industrial speaker",company="SKF",speakers="Lijun Cao",
      title="Integrating Particleworks into SKF Engineering Tools: Part I — Verification of Heat Transfer Coefficient",avatar="av_lijun-cao.png",sim="sim_skf.png"),
 dict(file="card-07-deepfluid.png",date=DATE,cat="Industrial speaker",company="DEEPFLUID",speakers="Dr. Lukas Hafner",
      title="Direct Optical Air-in-Oil Measurement for Gearings and Hydraulic Systems",avatar="av_lukas-hafner.png",sim="sim_deepfluid.png"),
 dict(file="card-08-trackone.png",date=DATE,cat="Industrial speaker",company="TRACK ONE",speakers="Leonardo Tiberi",
      title="Study of the Lubrication on Carrier Roller using a CFD 3D-MPS Method",avatar="av_leonardo-tiberi.png"),
 dict(file="card-09-univance.png",date=DATE,cat="Industrial speaker",company="UNIVANCE CORPORATION",speakers="Naohiro Fujita",
      title="Application of Particleworks to Gear Lubrication Analysis and Its Expansion to R&D on Airflow Effects",avatar="av_naohiro-fujita.png",sim="sim_univance.png"),
 dict(file="card-10-iav.png",date=DATE,cat="Industrial speaker",company="IAV",speakers="René Kockisch",
      title="From Formation to Dissolution: Air Bubble Dynamics in Gear Oil of an Electric 3-Speed Drivetrain",avatar="av_rene-kockisch.png"),
]

if __name__=="__main__":
    square_bg(os.path.join(PROC,"bg_square.png")); cal_icon(os.path.join(PROC,"cal.png"))
    for key,fn in PHOTOS.items():
        # original photo, no background removal (handled downstream by the user)
        make_avatar("photo",os.path.join(SPK,fn)).save(os.path.join(PROC,"av_"+key+".png"))
    make_avatar("ini","IM").save(os.path.join(PROC,"av_IM.png"))   # Issei Masaie — initials avatar
    make_sim(os.path.join(PROC,"skf.png"),os.path.join(PROC,"sim_skf.png"))
    make_sim(os.path.join(PROC,"deepfluid.png"),os.path.join(PROC,"sim_deepfluid.png"))
    make_sim(os.path.join(PROC,"univance.png"),os.path.join(PROC,"sim_univance.png"))
    print("assets ready")
    render_general()
    render_speakers()
    for c in cards: (render_keynote(c) if c["file"]=="card-01-prometech.png" else render_png(c))
    build_pptx(cards)
    print("ALL DONE")
